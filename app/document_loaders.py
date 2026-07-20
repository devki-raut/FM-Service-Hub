from collections import Counter
from pathlib import Path
from typing import Iterable
import re

import fitz
import pandas as pd
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".pdf", ".xlsx", ".xls", ".csv", ".pptx"}
EXCEL_ROWS_PER_RECORD = 100
EXCEL_SUMMARY_TOP_VALUES = 50
EXCEL_SUMMARY_TOP_TOKENS = 40


def iter_source_files(root: Path) -> Iterable[Path]:
    if root.is_file() and root.suffix.lower() in SUPPORTED_EXTENSIONS:
        yield root
        return

    for path in root.rglob("*"):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            yield path


def load_document(path: Path) -> list[dict]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _load_pdf(path)
    if suffix in {".xlsx", ".xls"}:
        return _load_excel(path)
    if suffix == ".csv":
        return _load_csv(path)
    if suffix == ".pptx":
        return _load_pptx(path)
    return []


def _load_pdf(path: Path) -> list[dict]:
    pages = []
    with fitz.open(path) as doc:
        for page_number, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                pages.append(
                    {
                        "text": text,
                        "metadata": {
                            "document_name": path.name,
                            "source_path": str(path),
                            "page": page_number,
                            "section": None,
                            "file_type": "pdf",
                        },
                    }
                )
    return pages


def _load_excel(path: Path) -> list[dict]:
    sheets = pd.read_excel(path, sheet_name=None, dtype=str)
    records = []
    for sheet_name, frame in sheets.items():
        frame = frame.dropna(how="all").fillna("")
        if frame.empty:
            continue

        records.extend(_excel_summary_records(path, sheet_name, frame))

        current_rows = []
        start_excel_row = 2
        for _, (index, row) in enumerate(frame.iterrows(), start=1):
            values = [f"{column}: {value}" for column, value in row.items() if str(value).strip()]
            if not values:
                continue

            if not current_rows:
                start_excel_row = int(index) + 2
            current_rows.append(f"Row {int(index) + 2}: " + "; ".join(values))

            if len(current_rows) >= EXCEL_ROWS_PER_RECORD:
                records.append(_excel_record(path, sheet_name, start_excel_row, int(index) + 2, current_rows))
                current_rows = []

        if current_rows:
            end_excel_row = start_excel_row + len(current_rows) - 1
            records.append(_excel_record(path, sheet_name, start_excel_row, end_excel_row, current_rows))
    return records


def _excel_summary_records(path: Path, sheet_name: str, frame: pd.DataFrame) -> list[dict]:
    records = []
    overview_lines = [
        f"Excel sheet summary for {path.name}, sheet {sheet_name}.",
        f"Total data rows: {len(frame)}.",
        f"Columns: {', '.join(str(column) for column in frame.columns)}.",
    ]
    unique_column = _first_column_containing(list(frame.columns), "complaint no")
    if unique_column:
        overview_lines.append(f"Unique {unique_column}: {frame[unique_column].nunique()}.")
    records.append(_excel_summary_record(path, sheet_name, "overview", overview_lines))

    date_lines = _excel_date_summary_lines(frame, unique_column)
    if date_lines:
        records.append(_excel_summary_record(path, sheet_name, "date-year counts", date_lines))

    categorical_lines = _excel_categorical_summary_lines(frame)
    if categorical_lines:
        records.append(_excel_summary_record(path, sheet_name, "column value counts", categorical_lines))

    token_lines = _excel_token_summary_lines(frame)
    if token_lines:
        records.append(_excel_summary_record(path, sheet_name, "text token counts", token_lines))

    return records


def _excel_summary_record(path: Path, sheet_name: str, summary_name: str, lines: list[str]) -> dict:
    return {
        "text": "\n".join(lines),
        "metadata": {
            "document_name": path.name,
            "source_path": str(path),
            "page": None,
            "section": f"{sheet_name} summary - {summary_name}",
            "file_type": "excel",
        },
    }


def _excel_date_summary_lines(frame: pd.DataFrame, unique_column: str | None) -> list[str]:
    lines = []
    date_columns = [column for column in frame.columns if "date" in str(column).lower()]
    for column in date_columns:
        dates = pd.to_datetime(frame[column], errors="coerce")
        years = sorted(dates.dt.year.dropna().astype(int).unique(), reverse=True)
        if not years:
            continue
        parts = []
        for year in years:
            year_frame = frame[dates.dt.year.eq(year)]
            part = f"{year}: {len(year_frame)} rows"
            if unique_column:
                part += f", {year_frame[unique_column].nunique()} unique {unique_column}"
            parts.append(part)
        lines.append(f"{column} year counts: {'; '.join(parts)}.")
    return lines


def _excel_categorical_summary_lines(frame: pd.DataFrame) -> list[str]:
    lines = []
    for column in frame.columns:
        values = frame[column].astype(str).str.strip()
        non_empty = values[values != ""]
        unique_count = non_empty.nunique()
        if unique_count < 2 or unique_count > EXCEL_SUMMARY_TOP_VALUES:
            continue
        if non_empty.str.len().mean() > 80:
            continue

        counts = non_empty.value_counts().head(EXCEL_SUMMARY_TOP_VALUES)
        rendered = "; ".join(f"{value}: {count}" for value, count in counts.items())
        lines.append(f"{column} value counts: {rendered}.")
    return lines


def _excel_token_summary_lines(frame: pd.DataFrame) -> list[str]:
    row_text = pd.Series([" ".join(row) for row in frame.astype(str).to_numpy()], index=frame.index)
    token_counts = Counter()
    for text in row_text:
        token_counts.update(set(_model_like_tokens(text)))

    common_tokens = [(token, count) for token, count in token_counts.most_common(EXCEL_SUMMARY_TOP_TOKENS) if count > 1]
    if not common_tokens:
        return []

    lines = ["Text token row counts for model/product/code-like values:"]
    created_column = _first_column_containing(list(frame.columns), "created date")
    created_dates = pd.to_datetime(frame[created_column], errors="coerce") if created_column else None

    for token, count in common_tokens:
        line = f"{token}: {count} rows"
        if created_dates is not None:
            matching_rows = row_text.str.contains(re.escape(token), case=False, na=False)
            years = sorted(created_dates[matching_rows].dt.year.dropna().astype(int).unique(), reverse=True)
            if years:
                year_parts = [f"{year}: {int((matching_rows & created_dates.dt.year.eq(year)).sum())} rows" for year in years]
                line += f"; {created_column} year counts: {'; '.join(year_parts)}"
        lines.append(line + ".")
    return lines


def _model_like_tokens(text: str) -> list[str]:
    tokens = re.findall(r"\b[A-Za-z]+[A-Za-z0-9-]*\d{2,}[A-Za-z0-9-]*\b", text)
    return [token.upper() for token in tokens if not token.isdigit() and len(token) <= 30]


def _first_column_containing(columns: list[str], needle: str) -> str | None:
    return next((column for column in columns if needle in str(column).lower()), None)


def _excel_record(path: Path, sheet_name: str, start_row: int, end_row: int, rows: list[str]) -> dict:
    return {
        "text": "\n".join(rows),
        "metadata": {
            "document_name": path.name,
            "source_path": str(path),
            "page": None,
            "section": f"{sheet_name} rows {start_row}-{end_row}",
            "file_type": "excel",
        },
    }


def _load_csv(path: Path) -> list[dict]:
    frame = pd.read_csv(path, dtype=str).dropna(how="all").fillna("")
    records = []
    current_rows = []
    start_csv_row = 2
    for index, row in frame.iterrows():
        values = [f"{column}: {value}" for column, value in row.items() if str(value).strip()]
        if not values:
            continue
        if not current_rows:
            start_csv_row = int(index) + 2
        current_rows.append(f"Row {int(index) + 2}: " + "; ".join(values))
        if len(current_rows) >= EXCEL_ROWS_PER_RECORD:
            records.append(
                {
                    "text": "\n".join(current_rows),
                    "metadata": {
                        "document_name": path.name,
                        "source_path": str(path),
                        "page": None,
                        "section": f"rows {start_csv_row}-{int(index) + 2}",
                        "file_type": "csv",
                    },
                }
            )
            current_rows = []
    if current_rows:
        records.append(
            {
                "text": "\n".join(current_rows),
                "metadata": {
                    "document_name": path.name,
                    "source_path": str(path),
                    "page": None,
                    "section": f"rows {start_csv_row}-{start_csv_row + len(current_rows) - 1}",
                    "file_type": "csv",
                },
            }
        )
    return records


def _load_pptx(path: Path) -> list[dict]:
    presentation = Presentation(path)
    slides = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(
                {
                    "text": "\n".join(texts),
                    "metadata": {
                        "document_name": path.name,
                        "source_path": str(path),
                        "page": slide_number,
                        "section": f"Slide {slide_number}",
                        "file_type": "pptx",
                    },
                }
            )
    return slides

